import os
import pandas as pd
from docx import Document
from pptx import Presentation

def get_questions(file_path, level):
    df = pd.read_json(file_path, lines=True)
    df = df[df["Level"] == level]
    
    result=[]
    
    for index, row in df.iterrows():
        result.append([row["Level"], row["Question"], row["file_name"], row["Final answer"]])

    return result

def is_ext(file_path, ext):
    return os.path.splitext(file_path)[1].lower() == ext.lower()
    
def read_file_json(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    df = None

    if ext == ".csv":
        df = pd.read_csv(file_path)
    elif ext in (".xls", ".xlsx"):
        df = pd.read_excel(file_path)
    elif ext in (".json", ".jsonl"):
        df = pd.read_json(file_path)

    return "" if df is None else df.to_json()

def read_docx_text(file_path):
    doc = Document(file_path)
    
    text = []

    for block in doc.element.body:
        if block.tag.endswith("p"):
            for paragraph in doc.paragraphs:
                if paragraph._element == block:
                    if paragraph.style.name.startswith("Heading"):
                        text.append("\n**" + paragraph.text + "**\n")
                    elif paragraph.text:
                        text.append(paragraph.text)
        elif block.tag.endswith("tbl"):
            for table in doc.tables:
                if table._element == block:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        text.append(" | ".join(row_text))
            
    return "\n".join(text)

def read_pptx_text(file_path):
    prs = Presentation(file_path)
    
    text = []
    
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text.append("\n".join(slide_text))
    
    return "\n\n".join(text)
